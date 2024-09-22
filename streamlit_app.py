import streamlit as st
from langchain_openai.chat_models import ChatOpenAI
import os
from dotenv import load_dotenv
import sqlite3
import datetime
import boto3
from botocore.exceptions import NoCredentialsError
import uuid
import pytesseract
from PIL import Image
import io
from PyPDF2 import PdfReader
from pdf2image import convert_from_bytes
from pptx import Presentation
import re
import whisper
import tempfile
from functools import lru_cache
from ai_ta import AI_TA
import streamlit_authenticator as stauth
import yaml
from yaml.loader import SafeLoader
import openai
from llama_index.core import SimpleDirectoryReader, VectorStoreIndex
# from llama_index.core.multi_modal_llms import OpenAIMultiModal

load_dotenv()

# Database setup
@st.cache_resource
def init_db():
    conn = sqlite3.connect('class_notes.db', check_same_thread=False)
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS classes
                 (id INTEGER PRIMARY KEY, 
                  name TEXT UNIQUE)''')
    
    c.execute('''CREATE TABLE IF NOT EXISTS ta_indexes
                 (id INTEGER PRIMARY KEY, 
                  class_name TEXT UNIQUE,
                  serialized_index TEXT UNIQUE,
                  FOREIGN KEY (class_name) REFERENCES classes(name))''')
    
    c.execute('''CREATE TABLE IF NOT EXISTS notes
                 (id INTEGER PRIMARY KEY, class_id INTEGER, content TEXT, timestamp DATETIME, file_urls TEXT, audio_urls TEXT, upvotes INTEGER DEFAULT 0, downvotes INTEGER DEFAULT 0,
                 FOREIGN KEY (class_id) REFERENCES classes(id))''')
    
    # Check if file_urls and audio_urls columns exist, if not, add them
    c.execute("PRAGMA table_info(notes)")
    columns = [column[1] for column in c.fetchall()]
    if 'file_urls' not in columns:
        c.execute("ALTER TABLE notes ADD COLUMN file_urls TEXT")
    if 'audio_urls' not in columns:
        c.execute("ALTER TABLE notes ADD COLUMN audio_urls TEXT")
    
    # Check if upvotes and downvotes columns exist, if not, add them
    if 'upvotes' not in columns:
        c.execute("ALTER TABLE notes ADD COLUMN upvotes INTEGER DEFAULT 0")
    if 'downvotes' not in columns:
        c.execute("ALTER TABLE notes ADD COLUMN downvotes INTEGER DEFAULT 0")
    
    conn.commit()
    return conn

if 'db_conn' not in st.session_state:
    st.session_state.db_conn = init_db()

# Helper functions
def add_class(class_name):
    c = st.session_state.db_conn.cursor()
    c.execute("INSERT OR IGNORE INTO classes (name) VALUES (?)", (class_name,))

    class_ta = AI_TA(class_name)
    c.execute("INSERT OR IGNORE INTO ta_indexes (class_name, serialized_index) VALUES (?, ?)", 
                (class_name, class_ta.serialize()))
    st.session_state.db_conn.commit()

def get_classes():
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT name FROM classes")
    return [row[0] for row in c.fetchall()]

def extract_text_from_file(file):
    text = ""
    if file.type == "application/pdf":
        pdf_reader = PdfReader(file)
        for page_num, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            if page_text.strip():
                text += page_text + "\n"
            else:
                # If no text was extracted, the page might be scanned. Use OCR.
                file.seek(0)  # Reset file pointer to the beginning
                images = convert_from_bytes(file.read(), first_page=page_num+1, last_page=page_num+1)
                for image in images:
                    text += pytesseract.image_to_string(image) + "\n"
    elif file.type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
        prs = Presentation(io.BytesIO(file.read()))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
            text += "\n"
    else:
        # Process as image
        image = Image.open(file)
        text = pytesseract.image_to_string(image)
    
    # Strip multiple consecutive new lines to be one new line at most
    text = re.sub(r'\n{2,}', '\n', text)
    return text.strip()

def process_files(files):
    documents = []
    for file in files:
        if file.type.startswith('image/'):
            img_data = file.read()
            documents.append(Document(image=img_data))
        elif file.type == 'application/pdf':
            pdf_reader = SimpleDirectoryReader().load_data(file)
            documents.extend(pdf_reader)
        # Add more file type handlers as needed
    return documents

def add_note(class_name, note_content, files=None, audio_files=None):
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT id FROM classes WHERE name = ?", (class_name,))
    class_id = c.fetchone()[0]
    
    file_urls = []
    audio_urls = []
    
    if files:
        for file in files:
            file_extension = os.path.splitext(file.name)[1]
            file_filename = f"files/{uuid.uuid4()}{file_extension}"
            if upload_to_s3(file, os.getenv('S3_BUCKET_NAME'), file_filename):
                file_urls.append(get_s3_url(os.getenv('S3_BUCKET_NAME'), file_filename))
    
    if audio_files:
        for audio_file in audio_files:
            audio_filename = f"audio/{uuid.uuid4()}{os.path.splitext(audio_file.name)[1]}"
            if upload_to_s3(audio_file, os.getenv('S3_BUCKET_NAME'), audio_filename):
                audio_urls.append(get_s3_url(os.getenv('S3_BUCKET_NAME'), audio_filename))
    
    # Strip multiple consecutive new lines in note_content
    note_content = re.sub(r'\n{2,}', '\n', note_content)
    
    c.execute("""INSERT INTO notes (class_id, content, timestamp, file_urls, audio_urls, upvotes, downvotes) 
                 VALUES (?, ?, ?, ?, ?, 0, 0)""",
              (class_id, note_content, datetime.datetime.now(), ','.join(file_urls), ','.join(audio_urls)))
    
    # Update the TA index
    c.execute("SELECT serialized_index FROM ta_indexes WHERE class_name = ?", (class_name,))
    serialized_index = c.fetchone()[0]
    class_ta = AI_TA.deserialize(serialized_index)
    print("class_ta.class_name: ", class_ta.class_name)
    print("class_ta.history: ", class_ta.train_history)
    class_ta.train(note_content, files)
    c.execute("UPDATE ta_indexes SET serialized_index = ? WHERE class_name = ?", (class_ta.serialize(), class_name))
    st.session_state.db_conn.commit()

def update_note(note_id, content, files=None, audio_files=None):
    c = st.session_state.db_conn.cursor()
    
    file_urls = []
    audio_urls = []
    
    if files:
        for file in files:
            file_extension = os.path.splitext(file.name)[1]
            file_filename = f"files/{uuid.uuid4()}{file_extension}"
            if upload_to_s3(file, os.getenv('S3_BUCKET_NAME'), file_filename):
                file_urls.append(get_s3_url(os.getenv('S3_BUCKET_NAME'), file_filename))
    
    if audio_files:
        for audio_file in audio_files:
            audio_filename = f"audio/{uuid.uuid4()}{os.path.splitext(audio_file.name)[1]}"
            if upload_to_s3(audio_file, os.getenv('S3_BUCKET_NAME'), audio_filename):
                audio_urls.append(get_s3_url(os.getenv('S3_BUCKET_NAME'), audio_filename))
    
    # Get existing file_urls and audio_urls
    c.execute("SELECT file_urls, audio_urls FROM notes WHERE id = ?", (note_id,))
    existing_file_urls, existing_audio_urls = c.fetchone()
    if existing_file_urls:
        file_urls = existing_file_urls.split(',') + file_urls
    if existing_audio_urls:
        audio_urls = existing_audio_urls.split(',') + audio_urls
    
    # Strip multiple consecutive new lines in content
    content = re.sub(r'\n{2,}', '\n', content)
    
    c.execute("""UPDATE notes 
                 SET content = ?, timestamp = ?, file_urls = ?, audio_urls = ?
                 WHERE id = ?""",
              (content, datetime.datetime.now(), ','.join(file_urls), ','.join(audio_urls), note_id))
    st.session_state.db_conn.commit()

def delete_file_from_note(note_id, file_url):
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT file_urls FROM notes WHERE id = ?", (note_id,))
    file_urls = c.fetchone()[0].split(',')
    file_urls.remove(file_url)
    c.execute("UPDATE notes SET file_urls = ? WHERE id = ?", (','.join(file_urls), note_id))

def get_notes(class_name):
    c = st.session_state.db_conn.cursor()
    try:
        c.execute("""SELECT notes.id, notes.content, notes.timestamp
                     FROM notes 
                     JOIN classes ON notes.class_id = classes.id 
                     WHERE classes.name = ?
                     ORDER BY notes.timestamp DESC""", (class_name,))
    except sqlite3.OperationalError:
        # If the query fails, fall back to the original schema
        c.execute("""SELECT notes.id, notes.content, notes.timestamp
                     FROM notes 
                     JOIN classes ON notes.class_id = classes.id 
                     WHERE classes.name = ?
                     ORDER BY notes.timestamp DESC""", (class_name,))
    return c.fetchall()

def get_tables():
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT name FROM sqlite_master WHERE type='table';")
    tables = c.fetchall()
    return [table[0] for table in tables]

# Add these new functions
def get_classes_data():
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT * FROM classes")
    return c.fetchall()

def get_notes_data():
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT * FROM notes")
    return c.fetchall()

# Add this new function
def clear_database():
    c = st.session_state.db_conn.cursor()
    c.execute("DELETE FROM notes")
    c.execute("DELETE FROM classes")
    # c.execute("DELETE FROM sqlite_sequence WHERE name IN ('notes', 'classes')")  # Reset auto-increment
    st.session_state.db_conn.commit()


# Add these new functions

def get_note_by_id(note_id):
    c = st.session_state.db_conn.cursor()
    c.execute("SELECT content FROM notes WHERE id = ?", (note_id,))
    return c.fetchone()[0]

# Add this new function to delete a note
def delete_note(note_id):
    c = st.session_state.db_conn.cursor()
    c.execute("DELETE FROM notes WHERE id = ?", (note_id,))
    st.session_state.db_conn.commit()

# Add these new functions for S3 operations
def upload_to_s3(file, bucket_name, object_name):
    pass
    s3_client = boto3.client('s3',
                             aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
                             aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY'))
    try:
        s3_client.upload_fileobj(file, bucket_name, object_name)
        return True
    except NoCredentialsError:
        st.error("AWS credentials not available")
        return False
    except Exception as e:
        st.error(f"An error occurred: {e}")
        return False

def get_s3_url(bucket_name, object_name):
    return f"https://{bucket_name}.s3.amazonaws.com/{object_name}"

# Add this cached function for audio transcription
@lru_cache(maxsize=None)
def cached_transcribe_audio(file_content):
    model = whisper.load_model("base")
    with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp_file:
        tmp_file.write(file_content)
        tmp_file_path = tmp_file.name

    result = model.transcribe(tmp_file_path)
    os.unlink(tmp_file_path)
    return result["text"]

def upvote_note(note_id):
    c = st.session_state.db_conn.cursor()
    c.execute("UPDATE notes SET upvotes = upvotes + 1 WHERE id = ?", (note_id,))
    st.session_state.db_conn.commit()

def downvote_note(note_id):
    c = st.session_state.db_conn.cursor()
    c.execute("UPDATE notes SET upvotes = MAX(upvotes - 1, 0) WHERE id = ?", (note_id,))
    st.session_state.db_conn.commit()

# Streamlit UI

# Load configuration file
with open('config.yaml') as file:
    config = yaml.load(file, Loader=SafeLoader)

# Create an authentication object
authenticator = stauth.Authenticate(
    config['credentials'],
    config['cookie']['name'],
    config['cookie']['key'],
    config['cookie']['expiry_days'],
    config['preauthorized']
)

# Add login widget
name, authentication_status, username = authenticator.login()

# Check authentication status
if authentication_status:
    authenticator.logout('Logout', 'sidebar')
    st.write(f'Welcome *{name}*')
    
    st.title("🎓 Class Notes Manager")

    # Add this new button for clearing the database
    # if st.button("Clear Database"):
    #     clear_database()
    #     st.success("Database cleared successfully!")
    #     st.experimental_rerun()

    # Add this new button
    # if st.button("Show Database Tables"):
    #     tables = get_tables()
    #     st.write("Tables in the database:", tables)

    # # Add these new buttons
    # if st.button("Show Classes Data"):
    #     classes_data = get_classes_data()
    #     st.write("Classes table contents:")
    #     st.table(classes_data)

    # if st.button("Show Notes Data"):
    #     notes_data = get_notes_data()
    #     st.write("Notes table contents:")
    #     st.table(notes_data)

    # Sidebar for class management
    st.sidebar.title("Class Management")

    # Display existing classes
    st.sidebar.subheader("Existing Classes")
    classes = get_classes()
    selected_class = st.sidebar.selectbox("Select a class:", [""] + classes)

    # Add new class
    new_class = st.sidebar.text_input("Add a new class:")
    if st.sidebar.button("Add Class"):
        if new_class:
            add_class(new_class)
            st.sidebar.success(f"Added {new_class}")
            st.rerun()
    # OpenAI API Key input
    openai_api_key = st.sidebar.text_input("OpenAI API Key", type="password", value=os.getenv("OPENAI_API_KEY", ""))

    # Main content area
    if selected_class:
        # AI-TA chat
        st.subheader("AI TA Chat :)")
        if openai_api_key.startswith("sk-"):
            chat_input = st.text_input("Ask AI-TA a question about this class:")
            if st.button("Ask"):
                c = st.session_state.db_conn.cursor()
                c.execute("SELECT serialized_index FROM ta_indexes WHERE class_name = ?", (selected_class,))
                serialized_index = c.fetchone()[0]
                print("serialized_index fetched! ")
                class_ta = AI_TA.deserialize(serialized_index)
                print("class_ta deseralized ", class_ta.class_name)
                openai.api_key = openai_api_key
                response = class_ta.query(chat_input)
                print("response generated ", response)
                st.info(response)

        st.subheader(f"Notes for {selected_class}")
        
        # Initialize session state for current note
        if 'current_note_id' not in st.session_state:
            st.session_state.current_note_id = None
        if 'current_note_content' not in st.session_state:
            st.session_state.current_note_content = ""


        # Text input for notes
        note_title = st.text_input("Note Title:")
        note_content = st.text_area("Edit note:", value=st.session_state.current_note_content)
        
        # File upload for multiple images/PDFs/PPTs
        files = st.file_uploader("Upload images, PDFs, or PowerPoint files (text extraction will be performed)", 
                                type=["png", "jpg", "jpeg", "pdf", "ppt", "pptx"], 
                                accept_multiple_files=True)
        if files:
            for i, file in enumerate(files):
                if file.type == "application/pdf":
                    st.write(f"PDF uploaded: {file.name}")
                elif file.type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
                    st.write(f"PowerPoint uploaded: {file.name}")
                else:
                    st.image(file, caption=f"Uploaded Image: {file.name}", width=200)
                extracted_text = extract_text_from_file(file)
                st.text_area(f"Extracted Text from {file.name}", value=extracted_text, height=150)
                if st.button(f"Append Extracted Text to Note (File {i+1})", key=f"append_button_{i}"):
                    note_content += f"\n\nExtracted Text from {file.name}:\n{extracted_text}"
                    st.session_state.current_note_content = note_content
                    st.rerun()
        
        # Audio file upload
        audio_files = st.file_uploader("Upload audio files (transcription will be performed)", 
                                    type=["mp3", "wav", "m4a"], 
                                    accept_multiple_files=True)
        
        if 'transcribed_texts' not in st.session_state:
            st.session_state.transcribed_texts = {}

        if audio_files:
            for i, audio_file in enumerate(audio_files):
                st.audio(audio_file)
                
                # Check if transcription is already done
                if audio_file.name not in st.session_state.transcribed_texts:
                    with st.spinner(f"Transcribing {audio_file.name}..."):
                        transcribed_text = cached_transcribe_audio(audio_file.read())
                    st.session_state.transcribed_texts[audio_file.name] = transcribed_text
                else:
                    transcribed_text = st.session_state.transcribed_texts[audio_file.name]
                
                st.text_area(f"Transcribed Text from {audio_file.name}", value=transcribed_text, height=150)
                if st.button(f"Append Transcribed Text to Note (Audio {i+1})", key=f"append_audio_button_{i}"):
                    note_content += f"\n\nTranscribed Text from {audio_file.name}:\n{transcribed_text}"
                    st.session_state.current_note_content = note_content
                    st.rerun()
        
        # Save button for the current note
        if st.button("Save Note"):
            if st.session_state.current_note_id:
                update_note(st.session_state.current_note_id, note_content, files, audio_files)
                st.success("Note updated!")
            else:
                add_note(selected_class, note_content, files, audio_files)
                st.success("New note added!")
            st.session_state.current_note_id = None
            st.session_state.current_note_content = ""
            st.rerun()
        
        # Display existing notes
        st.subheader("Existing Notes")

        c = st.session_state.db_conn.cursor()
        c.execute("""SELECT notes.id, notes.content, notes.upvotes
                     FROM notes 
                     JOIN classes ON notes.class_id = classes.id 
                     WHERE classes.name = ?
                     ORDER BY notes.upvotes DESC, notes.timestamp DESC""", (selected_class,))
        notes_data = c.fetchall()
        
        for note_id, note_content, upvotes in notes_data:
            container = st.container(border=True)
            
            # Top row with note ID, voting buttons, and counts
            top_col1, top_col2, top_col3, top_col4 = container.columns([3, 1, 1, 1])
            with top_col1:
                top_col1.markdown(f"### Note {note_id}")
            with top_col2:
                if st.button("👍", key=f"upvote_{note_id}"):
                    upvote_note(note_id)
                    st.rerun()
            with top_col3:
                st.write(f"Votes: {upvotes}")
            with top_col4:
                if st.button("👎", key=f"downvote_{note_id}"):
                    downvote_note(note_id)
                    st.rerun()
            
            # Content and action buttons
            content_col, action_col = container.columns([4, 1])
            with content_col:
                content_col.write(f"{note_content[:300]}...")
                image_number = note_id % 2 + 1
                content_col.image(f"user_profiles/{image_number}.jpg")
            with action_col:
                if st.button("Edit Note", key=f"edit_note_{note_id}"):
                    if st.session_state.current_note_id:
                        update_note(st.session_state.current_note_id, note_content)
                    st.session_state.current_note_id = note_id
                    st.session_state.current_note_content = get_note_by_id(note_id)
                    st.rerun()
                if st.button("Delete Note", key=f"delete_note_{note_id}"):
                    delete_note(note_id)
                    if st.session_state.current_note_id == note_id:
                        st.session_state.current_note_id = None
                        st.session_state.current_note_content = ""
                    st.rerun()

    # else:
    #     st.warning("Please enter your OpenAI API key to use the AI-TA feature!", icon="⚠")
else:
    st.info("Please select a class from the sidebar or add a new one.")

# Close the database connection when the app is done
# st.on_script_run.add(st.session_state.db_conn.close)
